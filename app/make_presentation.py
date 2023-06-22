import json
import re
from termcolor import colored as term_colored
from pptx import Presentation
from pptx.util import Inches


class PresentationCreator:
    def __init__(self):
        self.input_file = "../input/input.json"
        self.output_file = "../output/output.pptx"
        self.validation_max_content_words = 45
        self.validation_max_bullet_points = 8
        self.validation_max_title_length = 63

    def load_slide_data(self, filename):
        with open(filename, "r") as f:
            data = json.load(f)
        return data

    def refactor_data(self, data):
        if data["slides"][1]["type"] == "text":
            data["slides"][1]["content"] = data["slides"][1]["content"].replace("\n", "\n\n")
            data["slides"][1]["font_size"] = 13
        return data

    def validate_slide_data(self, data):
        slides = data["slides"]

        if slides[0]["type"] != "title" or slides[-1]["type"] != "title":
            raise ValueError("Presentation must start and end with a title slide")
        print(term_colored("Rule 1: Presentation starts and ends with a title slide - VALID", "green"))

        for slide in slides:
            if slide["type"] == "text":
                content = slide["content"]
                word_count = len(re.findall(r'\w+', content))
                if word_count > self.validation_max_content_words:
                    raise ValueError(
                        f"Slide '{slide['title']}' has more than {self.validation_max_content_words} words ({word_count})")
        print(term_colored("Rule 2: Slides do not exceed maximum word count - VALID", "green"))

        for slide in slides:
            if slide["type"] == "bullet_points":
                if len(slide["content"]) > self.validation_max_bullet_points:
                    raise ValueError(
                        f"Slide '{slide['title']}' has more than {self.validation_max_bullet_points} bullet points ({len(slide['content'])})")
        print(term_colored("Rule 3: Slides do not exceed maximum bullet points - VALID", "green"))

        for slide in slides:
            if len(slide['title']) > self.validation_max_title_length:
                raise ValueError(
                    f"Slide title '{slide['title']}' is longer than {self.validation_max_title_length} characters ({len(slide['title'])})")
        print(term_colored("Rule 4: Slide titles do not exceed maximum length - VALID", "green"))

    def create_title_slide(self, presentation, title):
        slide_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title

    def create_bullet_points_slide(self, presentation, title, content):
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        body_shape = slide.shapes.placeholders[1].text_frame
        for point in content:
            p = body_shape.add_paragraph()
            p.text = point
            p.level = 0

    def create_text_slide(self, presentation, title, content):
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        body_shape = slide.shapes.placeholders[1].text_frame
        body_shape.text = content

    def create_image_slide(self, presentation, title, image_path):
        slide_layout = presentation.slide_layouts[5]
        slide = presentation.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        left = Inches(1)
        top = Inches(1.5)
        pic = slide.shapes.add_picture(image_path, left, top, width=Inches(7.5))

    def generate_presentation(self):
        data = self.load_slide_data(self.input_file)
        data = self.refactor_data(data)
        self.validate_slide_data(data)
        slides = data["slides"]

        presentation = Presentation()
        for slide_data in slides:
            slide_type = slide_data["type"]
            if slide_type == "title":
                self.create_title_slide(presentation, slide_data["title"])
            elif slide_type == "bullet_points":
                self.create_bullet_points_slide(presentation, slide_data["title"], slide_data["content"])
            elif slide_type == "text":
                self.create_text_slide(presentation, slide_data["title"], slide_data["content"])
            elif slide_type == "image":
                self.create_image_slide(presentation, slide_data["title"], slide_data["image_path"])

        presentation.save(self.output_file)
        print('Presentation saved to ' + self.output_file)


if __name__ == "__main__":
    creator = PresentationCreator()
    creator.generate_presentation()
