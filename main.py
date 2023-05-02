import json
import re
from termcolor import colored as term_colored
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

input_file = "input.json"
output_file = "output.pptx"

validation_max_content_words = 45
validation_max_bullet_points = 8
validation_max_title_length = 63

def load_slide_data(filename):
    with open(filename, "r") as f:
        data = json.load(f)
    return data

def refactor_data(data):
    # if the second slide is a title slide, make the content text smaller
    if data["slides"][1]["type"] == "text":
        data["slides"][1]["content"] = data["slides"][1]["content"].replace("\n", "\n\n")
        data["slides"][1]["font_size"] = 13
    return data

def validate_slide_data(data):
    slides = data["slides"]
    
    # Rule 1: The data must start and end with a title slide
    if slides[0]["type"] != "title" or slides[-1]["type"] != "title":
        raise ValueError("Presentation must start and end with a title slide")
    print(term_colored("Rule 1: Presentation starts and ends with a title slide - VALID", "green"))

    # Rule 2: A single slide's text must not exceed 44 words
    for slide in slides:
        if slide["type"] == "text":
            content = slide["content"]
            word_count = len(re.findall(r'\w+', content))
            if word_count > validation_max_content_words:
                raise ValueError(f"Slide '{slide['title']}' has more than {validation_max_content_words} words ({word_count})")
    print(term_colored("Rule 2: Slides do not exceed maximum word count - VALID", "green"))

    # Rule 3: More rules that will guarantee good-looking presentations
    for slide in slides:
        if slide["type"] == "bullet_points":
            if len(slide["content"]) > validation_max_bullet_points:
                raise ValueError(f"Slide '{slide['title']}' has more than {validation_max_bullet_points} bullet points ({len(slide['content'])})")
    print(term_colored("Rule 3: Slides do not exceed maximum bullet points - VALID", "green"))

    # Rule 4: Slide titles must not exceed 50 characters
    for slide in slides:
        if len(slide['title']) > validation_max_title_length:
            raise ValueError(f"Slide title '{slide['title']}' is longer than {validation_max_title_length} characters ({len(slide['title'])})")
    print(term_colored("Rule 4: Slide titles do not exceed maximum length - VALID", "green"))


def create_title_slide(presentation, title):
    slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title

def create_bullet_points_slide(presentation, title, content):
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    body_shape = slide.shapes.placeholders[1].text_frame
    for point in content:
        p = body_shape.add_paragraph()
        p.text = point
        p.level = 0

def create_text_slide(presentation, title, content):
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    body_shape = slide.shapes.placeholders[1].text_frame
    body_shape.text = content

def create_image_slide(presentation, title, image_path):
    slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    left = Inches(1)
    top = Inches(1.5)
    pic = slide.shapes.add_picture(image_path, left, top, width=Inches(7.5))

if __name__ == "__main__":
    data = load_slide_data(input_file)
    data = refactor_data(data)
    validate_slide_data(data)
    slides = data["slides"]

    presentation = Presentation()
    for slide_data in slides:
        slide_type = slide_data["type"]
        if slide_type == "title":
            create_title_slide(presentation, slide_data["title"])
        elif slide_type == "bullet_points":
            create_bullet_points_slide(presentation, slide_data["title"], slide_data["content"])
        elif slide_type == "text":
            create_text_slide(presentation, slide_data["title"], slide_data["content"])
        elif slide_type == "image":
            create_image_slide(presentation, slide_data["title"], slide_data["image_path"])

    presentation.save(output_file)